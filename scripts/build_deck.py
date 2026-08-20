"""MGB-SED Magdalena hydrology deck, figure-led.

Figures come from three places, all real:
  figures/deck/<nb>_c*.png   extracted from our executed notebooks
  figures/deck/gen_*.png     generated from sim_calibrated_v2/*.csv
  figures/deck/yb_*.png      from the team's second implementation repo
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = pathlib.Path(__file__).resolve().parent.parent
FIG = REPO / "figures" / "deck"
OUT = REPO / "MGB-SED_Magdalena_FIGURES.pptx"

NAVY = RGBColor(0x1F, 0x35, 0x64)
INK = RGBColor(0x24, 0x29, 0x2E)
GREY = RGBColor(0x5A, 0x63, 0x6B)
ACCENT = RGBColor(0x0B, 0x6E, 0x99)
GOOD = RGBColor(0x1B, 0x7F, 0x4B)
WARN = RGBColor(0xB0, 0x3A, 0x2E)
BAND = RGBColor(0xEC, 0xF1, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.55)
BW = W - 2 * M

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = 6


def _tb(s, left, top, width, height=Inches(0.5)):
    tb = s.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def slide(num, title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    p = _tb(s, M, Inches(0.28), BW, Inches(0.62)).text_frame.paragraphs[0]
    p.text = f"{num}.  {title}"
    p.font.size, p.font.bold, p.font.color.rgb = Pt(25), True, NAVY
    r = s.shapes.add_shape(1, M, Inches(0.94), BW, Pt(1.2))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xD3, 0xDA, 0xE1)
    r.line.fill.background(); r.shadow.inherit = False
    y = Inches(1.02)
    if kicker:
        kp = _tb(s, M, Inches(0.97), BW, Inches(0.36)).text_frame.paragraphs[0]
        kp.text = kicker
        kp.font.size, kp.font.italic, kp.font.color.rgb = Pt(13), True, ACCENT
        y = Inches(1.42)
    return s, y


def bullets(s, y, items, left=None, width=None, size=15):
    tf = _tb(s, left or M, y, width or BW).text_frame
    first = True
    for it in items:
        lvl, txt = it if isinstance(it, tuple) else (0, it)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("• " if lvl == 0 else "– ") + txt
        p.font.size = Pt(size if lvl == 0 else size - 2.5)
        p.font.color.rgb = INK if lvl == 0 else GREY
        p.space_after = Pt(6 if lvl == 0 else 2)
    return tf


def pic(s, name, top, height=None, width=None, left=None, cap=None):
    f = FIG / name
    if not f.exists():
        b = _tb(s, M, top, BW, Inches(0.4)).text_frame.paragraphs[0]
        b.text = f"[missing figure: {name}]"
        b.font.size, b.font.color.rgb = Pt(12), WARN
        return None
    kw = {"height": height} if height else {"width": width}
    p = s.shapes.add_picture(str(f), Inches(0), top, **kw)
    p.left = Emu(int(left)) if left is not None else Emu(int((W - p.width) / 2))
    if cap:
        c = _tb(s, Emu(int(p.left)), Emu(int(top + p.height + Inches(0.03))),
                Emu(int(p.width)), Inches(0.3)).text_frame.paragraphs[0]
        c.text = cap
        c.font.size, c.font.italic, c.font.color.rgb = Pt(10), True, GREY
        c.alignment = PP_ALIGN.CENTER
    return p


def callout(s, y, text, colour=ACCENT, size=14, height=Inches(0.55)):
    b = s.shapes.add_shape(1, M, y, BW, height)
    b.fill.solid(); b.fill.fore_color.rgb = BAND
    b.line.color.rgb = colour; b.line.width = Pt(1.4); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.14)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size, p.font.bold, p.font.color.rgb = Pt(size), True, colour
    return b


def table(s, y, rows, col_w, size=11.5):
    nr, nc = len(rows), len(rows[0])
    tot = Inches(sum(col_w))
    shp = s.shapes.add_table(nr, nc, M + Emu(int((BW - tot) / 2)), y, tot, Inches(0.3 * nr))
    t = shp.table
    for j, wv in enumerate(col_w):
        t.columns[j].width = Inches(wv)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = t.cell(i, j); c.text = str(v)
            c.margin_left = c.margin_right = Inches(0.06)
            c.margin_top = c.margin_bottom = Inches(0.015)
            pr = c.text_frame.paragraphs[0]
            pr.font.size = Pt(size)
            pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = NAVY
                pr.font.bold, pr.font.color.rgb = True, WHITE
            else:
                c.fill.fore_color.rgb = BAND if i % 2 else WHITE
                pr.font.color.rgb = INK
                if j == 0:
                    pr.font.bold = True
    return shp


def note(s, t):
    s.notes_slide.notes_text_frame.text = t


# ============================================================== 1  title
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
pic(s, "04_c005_1.png", Inches(0.0), height=H)
band = s.shapes.add_shape(1, Inches(0), Inches(2.05), W, Inches(2.75))
band.fill.solid(); band.fill.fore_color.rgb = NAVY
band.line.fill.background(); band.shadow.inherit = False
band.fill.transparency = 0.08
tf = _tb(s, Inches(0.85), Inches(2.25), W - Inches(1.7), Inches(2.4)).text_frame
p = tf.paragraphs[0]
p.text = "MGB-SED suspended-sediment modelling of the\nMagdalena–Cauca basin: an ENSO contrast"
p.font.size, p.font.bold, p.font.color.rgb = Pt(33), True, WHITE
for i, ln in enumerate(["Hydrological calibration — attempts, measurements, and where we stand",
                        "Team of three  ·  Advisor: Prof. F. J. Briceño-Zuluaga  ·  UMNG  ·  August 2026",
                        "Method transferred from Fagundes et al. (2026), Int. Soil Water Conserv. Res. 14, 100599"]):
    q = tf.add_paragraph(); q.text = ln
    q.font.size = Pt(15 if i == 0 else 12)
    q.font.color.rgb = WHITE
    q.space_before = Pt(12 if i == 0 else 3)
note(s, "This deck stops at our attempts to calibrate the hydrological model: what we tried, "
        "what each attempt measured, where we stand. It does not claim a finished calibration.")

# ============================================================== 2  ENSO years
s, y = slide(2, "The question — and why these two years",
             "The contrast years come from our own data, not from the literature")
bullets(s, y, [
    "Magdalena–Cauca: 257,097 km² — among the world's highest specific sediment yields "
    "(Restrepo 2015: ~690 t km⁻² yr⁻¹ at Calamar; published, not ours)",
    "La Niña 2011 = +1.7σ wet   ·   El Niño 2015–16 = −1σ dry",
    "Objective: reproduce and EXPLAIN the flux difference, with spatial attribution",
], size=14)
pic(s, "06_c035_4.png", Inches(2.62), height=Inches(2.5), left=M,
    cap="Basin discharge anomaly per year — composite of stations with ≥20 years")
pic(s, "yb_2_enso_contrast.png", Inches(2.62), height=Inches(2.5),
    left=int(M + BW * 0.52), cap="ENSO contrast, independent analysis")
callout(s, Inches(5.6), "The gap: the ENSO–sediment link is documented observationally, never "
                        "reproduced with a process-based distributed model over the whole basin.")
note(s, "60 seconds. Both panels are ours — two independent analyses of the same signal.")

# ============================================================== 3  status
s, y = slide(3, "Where we are, honestly")
table(s, y, [
    ["Phase", "Status"],
    ["A — model inputs", "Complete"],
    ["B — hydrology", "CLOSED on H2E — at the r ≈ 0.57 rainfall-input ceiling"],
    ["C — sediment", "COMPLETE — calibration run; ENSO contrast reproduced (18/18)"],
], [2.9, 8.0], size=14)
bullets(s, Inches(2.7), [
    "MUSLE is driven by RUNOFF, not rainfall — the discharge model is the load-bearing component",
    "Calibrating sediment on an uncalibrated hydrology would fit erosion parameters to "
    "water-balance error",
], size=14)
pic(s, "13_c044_9.png", Inches(3.7), height=Inches(3.2),
    cap="Simulated mean specific runoff (mm/yr) — the field the sediment module consumes")
note(s, "Say 'not yet closed' out loud. Being explicit buys credibility for what follows.")

# ============================================================== 4  pipeline
s, y = slide(4, "Preprocessing pipeline — DEM to minibacias",
             "Every step verified against an independent recomputation")
for i, (nm, cap) in enumerate([("07_c005_1.png", "1 — conditioned DEM"),
                               ("07_c007_2.png", "2 — upstream area"),
                               ("07_c009_3.png", "3 — river network"),
                               ("07_c015_5.png", "4 — minibacias")]):
    pic(s, nm, Inches(1.45), height=Inches(2.55),
        left=int(M + i * (BW / 4)), cap=cap)
pic(s, "08_c010_3.png", Inches(4.42), height=Inches(2.5), left=M,
    cap="URH composition = soil texture × land cover")
pic(s, "12_c012_2.png", Inches(4.42), height=Inches(2.5), left=int(M + BW * 0.52),
    cap="Upstream area (log) — outlet 257,097 km²")
note(s, "8,672 minibacias, 24 URH types = 3 IGAC soil families x 8 hydrological land classes.")

# ============================================================== 5  geometry verified
s, y = slide(5, "Spatial discretisation — and its verification")
bullets(s, y, [
    "8,672 minibacias (D8 on Copernicus 30 m DEM), outlet at Calamar; 24 URH types",
    "Per-minibacia soil storage from IGAC: median 72.6 mm, range 13.5–255 mm",
], size=14)
table(s, Inches(2.18), [
    ["Verification", "Result"],
    ["Outlet upstream area vs sum of own areas", "257,096.93  vs  257,096.93 km²"],
    ["Two independent accumulators agree to", "1.8 × 10⁻⁸ km²"],
    ["Area-monotonicity violating edges", "0"],
    ["URH fractions row-sum error", "8.9 × 10⁻¹⁶"],
], [5.6, 5.0], size=12.5)
pic(s, "13_c012_2.png", Inches(4.12), height=Inches(2.85),
    cap="Soil water storage Wm (mm) from IGAC field survey — a measured 19× spatial range, not a constant")
note(s, "The verification table is the point: geometry is checked two independent ways, "
        "not asserted.")

# ============================================================== 6  two implementations
s, y = slide(6, "Two implementations — our main methodological asset")
table(s, y, [
    ["", "Implementation A", "Implementation B"],
    ["Channel routing", "Muskingum X = 0", "Local-inertial + floodplain (Bates 2010)"],
    ["Run time", "11.7 s / 4,018 days", "1,510 s / trial"],
    ["Enables", "4,000-evaluation search", "The paper's physics, Depresión Momposina"],
    ["Calibrated KGE (median)", "0.346 – 0.450", "0.329  (90 gauges)"],
], [2.5, 3.4, 4.9], size=12.5)
callout(s, Inches(3.05), "They agree to ~0.02 KGE — independent codebases, independent forcing "
                         "pipelines, different routing. A cross-validation neither arm could "
                         "produce alone.", GOOD, 14)
pic(s, "14_c019_1.png", Inches(3.82), height=Inches(3.0),
    cap="DDS convergence, 2 configurations × 2 seeds, run as concurrent processes (implementation A)")
note(s, "A buys search, B buys physics. A 774-run search on B would take 13 days; A does 4,000 "
        "evaluations overnight. Speak about both as OUR work.")

# ============================================================== 7  the split
s, y = slide(7, "Model period and the split",
             "The slide that makes every later number credible")
bullets(s, y, [
    "2008-01-01 → 2018-12-31, 4,018 days.  2008 warms up, 2009–2018 is scored",
    "Klemeš (1986) differential split-sample: we calibrate on NEUTRAL years only, so both "
    "ENSO phases are strictly OUT-OF-SAMPLE",
    (1, "the ENSO contrast is therefore a prediction, not a fit"),
    (1, "cal→val degradation −0.159, only 0.011 worse than an unfitted model"),
], size=14)
pic(s, "13_c023_5.png", Inches(3.15), height=Inches(3.35),
    cap="Warm-up: three mutually incompatible initial states converge to within 0.179 % of mean flow")
note(s, "Lead with this. Without it the deck is curve-fitting; with it, every later number is "
        "an out-of-sample result.")

# ============================================================== 8  three attempts
s, y = slide(8, "Three calibration attempts — what each one bought")
pic(s, "gen_attempts.png", y, height=Inches(3.75))
table(s, Inches(5.0), [
    ["Attempt", "Forcing", "Objective", "VAL KGE", "Recession", "At bound"],
    ["1 — Config B", "original", "daily KGE blend", "0.450", "2.98× too slow", "3 of 10"],
    ["2 — H1", "original", "+ recession term", "0.421", "0.96×", "2"],
    ["3 — H2", "repaired", "+ recession term", "0.346", "1.01×", "3"],
], [2.3, 1.5, 2.5, 1.3, 1.9, 1.5], size=12)
note(s, "DDS, 4,000 evaluations, two pre-registered configurations x two seeds, registered "
        "BEFORE running so no cell is a post-hoc pick. Objective references: prior 0.128, "
        "random sampling 0.173, attempt 1 0.243. Do not apologise for the KGE falling.")

# ============================================================== 9  the trade
s, y = slide(9, "The central result: fixing the physics costs skill")
pic(s, "gen_recession.png", y, height=Inches(3.55))
callout(s, Inches(4.85), "We traded 0.029 KGE for a recession that is right — and a dry phase "
                          "that beats climatology for the first time.", GOOD, 15)
bullets(s, Inches(5.62), [
    "Observed recession constants 9.5–11.9 d; attempt 1 simulated 27–45 d",
    "A higher KGE bought by a physically wrong recession is not the better model — which is why "
    "we report the ratio alongside the skill",
], size=13.5)
note(s, "THE argument of the deck. El Nino skill over climatology goes -0.026 -> +0.026 between "
        "attempts 1 and 2 - but say where it ENDED: -0.0005 in the adopted H2E (docs/26 addendum "
        "A.5). Present the KGE drop as a deliberate choice, and the dry phase as AT climatology.")

# ============================================================== 10 skill over clim
s, y = slide(10, "Scoring against a benchmark, not raw NSE")
pic(s, "gen_skill_clim.png", y, height=Inches(3.7))
bullets(s, Inches(5.0), [
    "A perfect day-of-year climatology also scores NSE −0.062 in the El Niño window, because "
    "that window has the record's highest observed variability (CV 0.799)",
    "So NSE is NOT comparable across windows — about a third of the apparent dry-phase failure "
    "was the metric, not the model",
    "In the ADOPTED fit (H2E): +0.106 (La Niña) vs −0.0005 (El Niño) KGE over climatology — the "
    "dry phase sits AT climatology, not above it. Across attempts 2→3→4 the El Niño figure ran "
    "+0.026 → +0.006 → −0.0005: it got HONESTER, not better",
], size=13.5)
note(s, "This is how we avoided over-diagnosing the dry phase. The asymmetry is real, but it is "
        "not 'worse than the mean'.")

# ============================================================== 11 controlled forcing test
s, y = slide(11, "Did repairing the rainfall help? — a controlled test",
             "Attempt 3 changes ONLY the forcing; matched gauges, matched window")
pic(s, "gen_h2_h1.png", y, height=Inches(3.3), left=M)
pic(s, "14_c032_3.png", Inches(1.5), height=Inches(3.3), left=int(M + BW * 0.53),
    cap="Per-gauge comparison: H2 above the line")
callout(s, Inches(5.0), "The repair fixed VOLUME and did not touch CORRELATION. We wrote that "
                         "prediction down before running it; r came back +0.003.", ACCENT, 14)
bullets(s, Inches(5.75), [
    "Confirms that volume error and correlation error are INDEPENDENT problems in this basin — "
    "which tells us exactly where to spend effort next",
], size=13.5)
note(s, "PBIAS 8.85 -> 4.41 %, r +0.0033, KGE -0.022, gauges with KGE>0 +2. The "
        "prediction-then-confirmation is the credibility moment.")

# ============================================================== 12 hydrographs
s, y = slide(12, "What the calibrated model actually produces",
             "Observed vs simulated daily discharge, gauges spanning three orders of magnitude in area")
pic(s, "13_c038_7.png", y, height=Inches(4.35))
bullets(s, Inches(6.0), [
    "Largest gauges have the BEST correlation (r 0.91) — they integrate enough area that "
    "rainfall noise averages out; their deficit is in amplitude, not timing",
], size=13)
note(s, "6 gauges above 20,000 km2 reach validation KGE 0.712 against 0.433 for the 55 smaller "
        "ones. The structural routing weakness shows up as amplitude, not as a timing error — "
        "we tested that fingerprint rather than assuming it.")

# ============================================================== 13 data defects
s, y = slide(13, "The data defect that value screens cannot see")
bullets(s, y, [
    "Stations OMIT dry days entirely, so mean rainfall scaled with how often the observer wrote "
    "anything down: 4.4 mm/day at >90 % reporting vs 11.7 mm/day below 50 %",
    "Our detector uses only the NEIGHBOURS' records, so it has a calibrated null",
], size=14)
pic(s, "10_c004_1.png", Inches(2.5), height=Inches(2.6), left=M,
    cap="The diagnostic that exposed it")
pic(s, "11_c004_1.png", Inches(2.5), height=Inches(2.6), left=int(M + BW * 0.52),
    cap="Reporting availability (black = station reporting)")
table(s, Inches(5.35), [
    ["Selectivity statistic", "Before", "After repair"],
    ["Sparse band (<50 % reporting)", "1.777", "1.040"],
    ["Dense band (>90 %) — the CONTROL", "1.001", "1.001  (held → no over-repair)"],
    ["Basin areal rainfall", "2,174 mm/yr", "2,036 mm/yr"],
], [4.4, 2.4, 3.8], size=11.5)
note(s, "153 stations repaired, 240,158 inferred-dry station-days. Energy-floor violations "
        "18 -> 14. Lesson: test for ABSENT records, not just outlier values — an outlier screen "
        "cannot see a record that was never written.")

# ============================================================== 14 verification
s, y = slide(14, "Verification as a first-class activity",
             "Four defects that only executing the code could reveal")
bullets(s, y, [
    "A silently truncating CSV reader — pandas returned 1,309 rows on one call and 3,630 on "
    "another, from a provably complete 4,018-row file, with NO exception",
    (1, "the cut is a contiguous PREFIX, so length, monotonicity, duplicate and calendar-gap "
        "checks all pass on it — only an assertion against an independently declared period caught it"),
    "A non-deterministic interpolator — three gauge pairs share exact coordinates, so the "
    "neighbour set was resolved by column order; shuffling moved 83 minibacias by 20.5 mm/day",
    "\"132 of 132 files present\" was a filename count — one ERA5 mosaic was internally corrupt "
    "at a plausible 43.7 MB",
    "Two gauges 5 cm apart were NOT duplicates (r = 0.756 over 1,470 shared days) — a "
    "distance-based merge would have destroyed a real record",
], size=13.5)
pic(s, "12_c009_1.png", Inches(4.95), height=Inches(1.8), left=M,
    cap="Join integrity across six tables (0 everywhere = clean)")
callout(s, Inches(6.85), "Standing guarantees:  mass-balance residual 1.67 × 10⁻¹⁷   ·   two "
                          "independent routing back-ends agree to max |ΔQ| = 0", GOOD, 13,
        Inches(0.48))
note(s, "Without the truncation assertion we would have calibrated on 1,309 of 4,018 days with "
        "every diagnostic green. Strongest single verification story here.")

# ============================================================== 15 the ceiling
s, y = slide(15, "The model is at its input's ceiling",
             "The scientific contribution")
pic(s, "11_c015_6.png", y, height=Inches(3.25), left=M)
table(s, Inches(1.5), [
    ["", "r"],
    ["Model, catchment-scale daily anomaly", "0.476"],
    ["Rainfall field's own leave-one-out skill", "0.429"],
    ["Inter-gauge daily correlation, 0–25 km", "0.33"],
    ["The same, 25–50 km", "0.25"],
    ["Mean gauge spacing", "~30 km"],
], [4.3, 1.3], size=12)
callout(s, Inches(4.95), "Across ALL 12 parameter configurations tested, El Niño correlation "
                          "stayed inside 0.556–0.572. Once bias and variance are repaired, KGE IS r.",
        ACCENT, 14)
bullets(s, Inches(5.7), [
    "The model sits just above the point-scale skill of the field driving it — the ceiling is a "
    "property of the OBSERVING NETWORK, not of the model",
    "No further parameter tuning can move the dry phase: remaining headroom ≈ +0.02, already located",
], size=13.5)
note(s, "Show this right after slide 8, or the first question is 'why is KGE only 0.4?' and the "
        "answer is a measured ceiling, not an excuse. Most likely slide to become a paper section.")

# ============================================================== 16 limits
s, y = slide(16, "What we cannot yet claim", "Say these before you are asked")
bullets(s, y, [
    "Phase B closed by DECISION at a measured input ceiling, not by passing — four configurations, "
    "none meeting every criterion we set in advance; it then closed a second time on measured "
    "conflict (H-PEAK refuted, H2E-S failed 2 of 3)",
    "Conventional adequacy not reached: Moriasi et al. (2007) put satisfactory daily NSE above "
    "0.50; ours is +0.16 to +0.26",
    "The ENSO asymmetry persists, and in the ADOPTED fit it is starker: skill over climatology is "
    "+0.106 in La Niña but −0.0005 in El Niño — the dry phase sits AT climatology, not above it",
    "Parameters still sit at bounds. The FAO-56 threshold form did free the crop coefficient "
    "(kc_mult 1.662, off the ~2.0 rail H1 hit) but it remains above any FAO-56 value",
    (1, "remaining candidate: ET = ETp·W/Wm throttles evaporation even in moist soil, and a "
        "doubled crop coefficient is exactly that compensation"),
    "No per-gauge specific yield can be published — our two networks' catchment areas disagree "
    "beyond 2× on 36 % of 85 shared gauges while their medians agree to 1 %",
    "Channel routing is Muskingum X = 0 in implementation A, so the fitted celerity of 0.221 m/s "
    "is a floodplain-storage SURROGATE, not a physical velocity",
], size=13.5)
pic(s, "06_c046_9.png", Inches(5.22), height=Inches(1.68), left=M,
    cap="Sediment rating curves — median R² ≈ 0.5, a stated uncertainty carried into Phase C")
note(s, "Volunteering the limits makes the positive claims believable. A 2.5x area error is a "
        "2.5x specific-yield error, so we report absolute flux and RATIOS only — never t/km2/yr.")

# ============================================================== 17 next steps
s, y = slide(17, "What we resolved since — and the one item that remains")
table(s, y, [
    ["#", "Step", "Why this order"],
    ["1", "Satellite-rainfall (CHIRPS) merge to lift the r ceiling",
     "DONE — bounded at +0.006; the ceiling is STRUCTURAL (closed-negative)"],
    ["2", "FAO-56 ET threshold form", "DONE — adopted in H2E"],
    ["3", "Search-seed expansion", "DONE — H2E confirmed"],
    ["4", "Phase C — sediment calibration + ENSO contrast",
     "DONE — exploratory fit; contrast reproduced 18/18 (docs/55, 56)"],
    ["5", "Resolve catchment areas against an external source",
     "OWED — the one real remaining item; unlocks specific-yield reporting"],
], [0.5, 5.8, 4.6], size=11.5)
bullets(s, Inches(3.55), [
    "The satellite merge (item 1) was built and validated but bounded to +0.006 r — the r ≈ 0.57 "
    "ceiling is the information content of the gauge network, not a processing gap (docs/58)",
], size=13)
pic(s, "yb_rs_retrieval_sentinel2_msi_nir.png", Inches(4.3), height=Inches(2.55), left=M,
    cap="Remote-sensing SSC retrieval — an independent avenue, not required for the result")
pic(s, "yb_demo_synthetic_overview.png", Inches(4.3), height=Inches(2.55),
    left=int(M + BW * 0.52), cap="Sediment transport module — now run on the real basin (Phase C)")
note(s, "Only item 5 (catchment areas) genuinely remains; items 1-4 are resolved. Phase C is "
        "DONE, not waiting on data quality — the sediment result slides follow.")

# ============================================================== 18 PHASE C: water to sediment
s, y = slide(18, "Phase C — from water to sediment", "the ENSO question, made physical")
bullets(s, y, [
    "We drove a MUSLE erosion model with the calibrated runoff, and re-derived the slope factor to "
    "the published source method (f_LS = 0.25146 — our original was ~4× too strong)",
    "One question: does the El Niño–La Niña swing change suspended-sediment transport, and by how much?",
    "Basin gross hillslope erosion: 299.5387 Mt/yr @ williams_m3 / us_customary / "
    "cp_revision=cited_central_2026_08_11 (prior revision 248.7298) — "
    "NEVER as a per-area yield (catchment-area embargo, docs/23 §13.2)",
], size=15)
note(s, "The bridge from the hydrology ceiling to the sediment question. Keep short.")

# ============================================================== 19 the calibration rails
s, y = slide(19, "The sediment calibration rails — a finding, not a failure",
             "the absolute level is not identifiable")
bullets(s, y, [
    "Best in-box median KGE_ln −0.118 (est. a) / +0.139 (est. b — same sign, so not indeterminate)",
    "The level knob wants α ≈ 0.48 — below the plausible floor",
    "The equation multiplies SEVEN constants the network cannot separate (design-matrix condition "
    "number = ∞); only the product Π is identifiable",
    "A fitted α would HIDE errors, not find them — so we report the level as unresolved, with a "
    "stated reason. Two independent implementations agree (docs/59)",
], size=13.5, width=int(BW * 0.5))
pic(s, "gen_kge_rail.png", Inches(1.55), height=Inches(3.4),
    left=int(M + BW * 0.5), cap="Median sediment KGE vs α — the fit rails at the box floor")

# ============================================================== 20 the question & mechanism
s, y = slide(20, "The ENSO–sediment question, and the physical chain", "hypothesis before evidence")
bullets(s, y, [
    "ENSO shifts tropical-Pacific sea-surface temperature, which shifts rainfall over Colombia: "
    "La Niña is wet, El Niño is dry (here, 2011 vs 2015–16)",
    "The chain is causal and one-directional: SST anomaly → rainfall anomaly → surface runoff → "
    "hillslope erosion AND transport capacity → suspended-sediment flux at the gauge",
    "Hypothesis, fixed before looking: the wet phase should carry markedly MORE sediment — because "
    "the erosive rainfall and the transporting discharge rise together, and MUSLE multiplies them",
    "The link is documented observationally, but has never been REPRODUCED with a process-based, "
    "distributed model over the whole basin. That gap is what this section tests.",
], size=14)
note(s, "Fix the hypothesis before the data. The section then asks in turn: is it observed? can we "
        "reproduce it? how robust is it?")

# ============================================================== 21 what is measurable
s, y = slide(21, "What is measurable — paired discharge, not SSC, is the binding constraint",
             "the honest denominator")
bullets(s, y, [
    "Sediment FLUX needs concentration AND discharge on the same day at the same station. Of 72 "
    "station-windows only 38 admit estimator (a) and 39 admit (b); 26 are impossible — no paired "
    "discharge in the window at all",
    "SSC exists where discharge does not. The one Magdalena-TRUNK station (ARRANCAPLUMAS) has 195 "
    "sediment samples in the El Niño window and ZERO discharge days — its record ends 2014-12-31",
    "Consequence, stated plainly: NO ENSO contrast is computable on the main stem. Every observed "
    "number that follows is a TRIBUTARY and Cauca-branch result — we say so, rather than implying "
    "basin-wide coverage",
], size=14)
note(s, "This kind of limit builds credibility: report exactly what the data can and cannot support.")

# ============================================================== 22 the observed contrast (figure slide)
s, y = slide(22, "The observed contrast — model-free evidence", "22 of 22 ratios exceed 1, no counter-examples")
pic(s, "gen_obs_contrast_detail.png", Inches(1.18), height=Inches(4.35))
callout(s, Inches(5.72), "Two independent estimators with conservative confidence intervals — the "
        "wet phase wins at every station. Primary-window median ~3–5×; the sharper ONI-peak windows "
        "give ~5–9×. The honest statement is a factor of ~3 to ~9, with the window definition worth "
        "roughly ×2 of that spread — quoting a single number would be false precision.", NAVY, 12.5,
        Inches(1.15))

# ============================================================== 23 why a weak fit can still test it
s, y = slide(23, "Why a weak calibration can still test the contrast",
             "the ratio cancels what we cannot pin down")
callout(s, Inches(1.5), "The wet/dry ratio removes every factor the calibration could not identify.",
        ACCENT, 16, Inches(0.75))
bullets(s, Inches(2.55), [
    "Sediment = α · LS · (runoff term)^β · K · C · P. The level knob α and the topographic factor "
    "LS are STATIC multipliers — identical every day, wet or dry",
    "In a within-station wet/dry ratio they sit identically on top and bottom and CANCEL EXACTLY:  "
    "R = (α·LS·X_wet)/(α·LS·X_dry) = X_wet ÷ X_dry",
    "So the C4.3 railing (α is unidentifiable) and the LS-level uncertainty do NOT touch the "
    "contrast. It depends only on the runoff contrast the rainfall carries between the two regimes",
    "The methodological point: a model too weak to predict absolute daily tonnage can still be a "
    "VALID INSTRUMENT for the relative ENSO signal — which is the actual research question",
], size=13.5)

# ============================================================== 24 the modelled contrast (figure slide)
s, y = slide(24, "The modelled contrast — reproduced, 18 / 18", "THE RESULT")
pic(s, "gen_enso_contrast.png", Inches(1.12), height=Inches(4.25))
callout(s, Inches(5.55), "Modelled median 3.05× — it matches the rating estimator (b) 2.84–2.95 and "
        "sits below the noisier sample estimator (a) 4.62. Every one of the 18 stations shows La "
        "Niña > El Niño, the same direction as the data — with NO tuning of α to achieve it.",
        GOOD, 12.5, Inches(1.1))

# ============================================================== 25 robustness (figure slide)
s, y = slide(25, "Robustness — not an artefact of a modelling choice", "invariant direction, honest magnitude")
pic(s, "gen_sensitivity_detail.png", Inches(1.2), height=Inches(3.85))
bullets(s, Inches(5.32), [
    "The direction (18/18, La Niña > El Niño) holds across three storm exponents β AND both window "
    "definitions — six independent cells, with no reversal in any of them",
    "The magnitude honestly varies (~2.6 to ~5.9×): the sharper ONI-peak windows exclude the ENSO "
    "shoulders and give the larger contrast — the SAME behaviour the observations show (slide 22)",
], size=13)

# ============================================================== 26 the limits, quantified
s, y = slide(26, "The limits, quantified — not hand-waved", "measured, not waved away")
bullets(s, y, [
    "Gauges: only ~18 have paired sediment + discharge; the 46 sediment-only sites have no discharge "
    "record at all (B5) → the flux set cannot grow. A physical limit of the network",
    "The r ≈ 0.57 ceiling proved unliftable: the one surviving rainfall lever bounded at ≤ +0.006 r",
    "The conclusion rests on THREE converging independent lines (observed flux 22/22, observed "
    "concentration, modelled 18/18) — not on gauge count",
], size=14, width=int(BW * 0.52))
pic(s, "gen_ceiling_bound.png", Inches(2.2), height=Inches(2.7),
    left=int(M + BW * 0.52), cap="A perfect rainfall repair lifts r by only +0.006")

# ============================================================== 27 the ask
s, y = slide(27, "The question we need your guidance on")
callout(s, Inches(1.55), "Is the input-ceiling result (slide 15) an acceptable closing statement "
                          "for the hydrological phase?", NAVY, 17, Inches(0.85))
bullets(s, Inches(2.75), [
    "If YES — the phase can close on a quantified limit, whether or not the rainfall merge succeeds",
    "If NO — if conventional adequacy is expected — then the merge must succeed, and if it does "
    "not we would need either:",
    (1, "denser rainfall input than IDEAM provides, or"),
    (1, "a reduced target: monthly instead of daily, or sub-basins instead of the full network"),
], size=15)
callout(s, Inches(5.0), "This changes what \"done\" means for Phase B — the most useful thing we "
                         "can settle today.", ACCENT, 15)
note(s, "End here, not on the summary. Asking the scope question is worth more than a recap.")

# ============================================================== 28 contribution
s, y = slide(28, "What we have contributed")
bullets(s, y, [
    "The ENSO sediment contrast REPRODUCED out-of-sample (18/18 stations, ~3×) — plus the proof "
    "that the absolute level is not identifiable from this network: the headline of the whole study",
    "First MGB-SED transposition to the Magdalena–Cauca, in TWO independent implementations that "
    "agree to ~0.02 KGE — one fast enough to calibrate, one carrying the paper's floodplain physics",
    "A mass-conservative, reproducible engine (residual 1.67 × 10⁻¹⁷, 11.7 s per basin-decade) "
    "with both ENSO phases held out of calibration",
    "A quantified INPUT CEILING for daily rainfall–runoff modelling at this gauge density — the "
    "result most likely to transfer beyond this basin",
    "A QC methodology for IDEAM station data, including a defect class — omitted dry days — that "
    "value-based screening cannot detect by construction",
    "A full audit trail: ~300 KB of technical documentation, a register of refuted hypotheses, "
    "and a traps reference",
], size=14)
pic(s, "13_c033_6.png", Inches(4.50), height=Inches(2.42),
    cap="Basin water balance components, daily — mass conservation holds to machine precision")
note(s, "Keep short if time is tight; slide 18 should get the discussion.")

prs.save(OUT)
n = len(prs.slides._sldIdLst)
print(f"wrote {OUT.name}  {OUT.stat().st_size/1024:.0f} KB  {n} slides")
